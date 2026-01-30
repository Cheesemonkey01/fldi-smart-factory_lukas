from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

prs = Presentation()
now = datetime.now().strftime('%Y-%m-%d')

def add_slide(title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            body.text = b
        else:
            p = body.add_paragraph()
            p.text = b
            p.level = 1
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Smart Factory Light Cell'
slide.placeholders[1].text = f'Präsentation\nAutor: Lukas\nDatum: {now}'
slide.notes_slide.notes_text_frame.text = 'Kurze Einführung: Projektziel und Demo-Plan.'

# Slides
add_slide('Motivation & Ziele', [
    'Vernetzte Smart-Factory-Zelle als CPS',
    'Technische Umsetzung, digitale Modellierung, Versionierung',
    'Lernziele: IoT, MQTT, HMI, Docker, Embedded Logik'
], notes='Erkläre kurz die Motivation und Lernziele.')

add_slide('Systemübersicht / Architektur', [
    'Komponenten: Sensor Simulator, Controller, Mosquitto Broker, Flutter HMI',
    'Kommunikation: MQTT Topics (Telemetry, Command, State)'
], notes='Zeige Architekturdiagramm (verbale Beschreibung).')

add_slide('Datenfluss & Topics', [
    'Telemetry: factory/lightcell/telemetry/temperature',
    'Commands: factory/lightcell/command/setpoint, /mode, /cooling',
    'State: factory/lightcell/state/cooling',
    'Payload-Beispiele: {"temperature": 25.3}, {"cooling_on": true}'
], notes='Kurz die wichtigsten Topics und Beispiel-Payloads vorlesen.')

add_slide('Python Backend', [
    'Sensor Simulator: erzeugt Temperaturwerte, reagiert auf cooling state',
    'Controller: verwaltet Setpoint, Mode, Cooling-Logic; empfängt Commands und publiziert State'
], notes='Erkläre evaluate_cooling() und den MANUAL-Modus.')

add_slide('Flutter HMI', [
    'Live-Temperaturanzeige, Setpoint-Input, Mode-Auswahl (AUTO/MANUAL)',
    'Manueller Schalter für Kühlung (sichtbar in MANUAL)',
    'MQTT-Integration mittels mqtt_client Package'
], notes='Screenshots zeigen falls verfügbar; ansonsten UI-Features erläutern.')

add_slide('Docker & Mosquitto', [
    'Mosquitto in Docker-Compose (Ports: 1883, 9001)',
    'Konfig: mosquitto.conf (persistence, websockets, allow_anonymous true für Demo)'
], notes='Erkläre kurz wie man den Broker startet: docker-compose up -d')

# Readme Hindernisse (summarise)
hindernisse = [
    'Firewall / Port-Probleme (Windows)',
    'OneDrive Sync kann Dateien sperren',
    'MQTT Payload-Inkompatibilitäten',
    'Reconnect-Strategien und Race-Conditions'
]
add_slide('Hindernisse & Lessons Learned', hindernisse, notes='Tipps zur Behebung kurz anreißen.')

add_slide('Tests & Smoke-Test', [
    'Smoke-Test: python/smoke_test.py prüft MANUAL/Cooling Verhalten',
    'Integrationstest: HMI + Controller + Broker + Simulator zusammen starten'
], notes='Erkläre Ablauf des Smoke-Tests und wie man ihn ausführt.')

add_slide('Demo-Ablauf (Live)', [
    '1) Start Docker (Mosquitto), Sensor-Simulator, Controller, HMI',
    '2) Zeige Mode-Wechsel AUTO → MANUAL',
    '3) Schalte manuelle Kühlung ein/aus und beobachte Sensor/State'
], notes='Kurzanleitung für Live-Demo, was Zeigen und welche Logs beobachten.')

add_slide('Nächste Schritte & Verbesserungen', [
    'TLS/Auth für Mosquitto (Produktionsreife)',
    'UI Verbesserungen und Hysterese in der Steuerlogik',
    'CI-Integration für Smoke-Tests'
], notes='Vorschläge für Weiterentwicklung nennen.')

add_slide('Fragen', ['Vielen Dank!'], notes='Bereit für Q&A.')

# Save
prs.save('docs/Presentation.pptx')
print('Presentation written to docs/Presentation.pptx')
